from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ── Brand palette ────────────────────────────────────────────────────────────
DARK_GREEN  = RGBColor(0x1A, 0x3C, 0x2E)
MID_GREEN   = RGBColor(0x2E, 0x7D, 0x5A)
LIGHT_GREEN = RGBColor(0x6F, 0xCF, 0x97)
ACCENT_GOLD = RGBColor(0xF2, 0xC9, 0x4E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xCC, 0xCC, 0xCC)
DARK_ROW    = RGBColor(0x1E, 0x44, 0x36)

BASE = "/Users/I767493/Library/CloudStorage/OneDrive-SAPSE/DHSN/2. Semester/Englisch"
CHARTS = os.path.join(BASE, "charts")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ── Helpers ───────────────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, colour=DARK_GREEN):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour

def rect(slide, left, top, width, height, fill_colour):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_colour
    return shape

def txbox(slide, text, left, top, width, height,
          size=24, bold=False, colour=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    run.font.name = "Calibri"
    return tb

def hline(slide, top, left=0.5, right=12.83, colour=MID_GREEN, thickness=Pt(2)):
    line = slide.shapes.add_connector(1, Inches(left), top, Inches(right), top)
    line.line.color.rgb = colour
    line.line.width = thickness

def label_tag(slide, text, left=Inches(0.5), top=Inches(6.9)):
    tb = slide.shapes.add_textbox(left, top, Inches(5), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = MID_GREEN
    run.font.name = "Calibri"
    run.font.bold = True

def add_bullets(slide, items, left, top, width, height, size=20):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.space_before = Pt(8)

def styled_table(slide, headers, rows, left, top, width, height):
    cols = len(headers)
    tbl = slide.shapes.add_table(len(rows)+1, cols, left, top, width, height).table
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid(); cell.fill.fore_color.rgb = MID_GREEN
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14); p.font.bold = True
        p.font.color.rgb = WHITE; p.font.name = "Calibri"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri+1, ci)
            cell.text = val
            bg_c = DARK_ROW if ri % 2 == 0 else RGBColor(0x22, 0x4D, 0x3E)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg_c
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE; p.font.name = "Calibri"

def png(slide, name, left, top, width, height):
    path = os.path.join(CHARTS, name + ".png")
    slide.shapes.add_picture(path, left, top, width, height)

def slide_header(slide, label, title, title_size=40):
    txbox(slide, label, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
          size=13, bold=True, colour=MID_GREEN)
    txbox(slide, title, Inches(0.5), Inches(0.8), Inches(12), Inches(0.9),
          size=title_size, bold=True)
    hline(slide, Inches(1.85))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), MID_GREEN)

# Decorative circles via XML transparency
for cx, cy, r, alpha_val in [(11.5, 5.5, 4.0, 0.06), (10.0, 1.5, 2.5, 0.08)]:
    shape = s.shapes.add_shape(9, Inches(cx-r), Inches(cy-r), Inches(r*2), Inches(r*2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MID_GREEN
    shape.line.fill.background()
    solidFill = shape.fill._xPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
            alpha_elem.set('val', str(int(alpha_val * 100000)))

txbox(s, "NexVerde", Inches(0.5), Inches(1.6), Inches(10), Inches(1.8),
      size=80, bold=True, colour=WHITE)
txbox(s, "Where performance meets the planet.",
      Inches(0.5), Inches(3.4), Inches(10), Inches(0.7),
      size=22, colour=LIGHT_GREEN)
hline(s, Inches(4.25), left=0.5, right=8.0)
txbox(s, "Sustainable Data Centre Proposal",
      Inches(0.5), Inches(4.4), Inches(9), Inches(0.5), size=17, colour=LIGHT_GREY)
txbox(s, "Commissioned for Birdy Capital",
      Inches(0.5), Inches(4.9), Inches(9), Inches(0.5), size=17, colour=LIGHT_GREY)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Who We Are
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Who We Are", "NexVerde — Built for this.")

cards = [
    ("12", "years in sustainable\ninfrastructure"),
    ("€200M+", "in delivered projects\nacross Europe"),
    ("ISO 50001", "certified · EU\nTaxonomy aligned"),
]
for i, (big, small) in enumerate(cards):
    x = Inches(0.5 + i * 4.2)
    rect(s, x, Inches(2.1), Inches(3.9), Inches(3.6), RGBColor(0x1E, 0x44, 0x36))
    txbox(s, big, x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(1.2),
          size=38, bold=True, colour=LIGHT_GREEN, align=PP_ALIGN.CENTER)
    txbox(s, small, x + Inches(0.2), Inches(3.7), Inches(3.5), Inches(1.0),
          size=17, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)
label_tag(s, "Speaker: Person A  ·  2 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Your Brief
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Your Brief", "What you asked for — and what we commit to.")

commits = [
    ("TIER III", "99.982% uptime\nFully redundant systems"),
    ("NET-ZERO", "Carbon trajectory\nMeasurable milestones"),
    ("€50M / 4 YRS", "On budget\nOn time delivery"),
]
for i, (tag, desc) in enumerate(commits):
    x = Inches(0.5 + i * 4.2)
    rect(s, x, Inches(2.1), Inches(3.9), Inches(4.0), RGBColor(0x22, 0x4D, 0x3E))
    rect(s, x, Inches(2.1), Inches(3.9), Inches(0.12), MID_GREEN)
    txbox(s, tag, x + Inches(0.15), Inches(2.35), Inches(3.6), Inches(0.8),
          size=26, bold=True, colour=LIGHT_GREEN, align=PP_ALIGN.CENTER)
    txbox(s, desc, x + Inches(0.15), Inches(3.2), Inches(3.6), Inches(1.5),
          size=17, colour=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Full brief in your private portal",
      Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
      size=14, colour=LIGHT_GREY, align=PP_ALIGN.CENTER)
label_tag(s, "Speaker: Person A  ·  3 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Technical Architecture  (PUE chart from PNG)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Technical Architecture", "The build.")

add_bullets(s, [
    "Modular design — scalable from day one",
    "Tier III redundancy: N+1 on all critical systems",
    "Liquid cooling core  ·  Air-side economisation",
    "Full schematics in your portal",
], Inches(0.5), Inches(2.1), Inches(6.5), Inches(3.8), size=20)

png(s, "chart_pue", Inches(7.0), Inches(1.85), Inches(5.9), Inches(4.6))
label_tag(s, "Speaker: Person B  ·  5 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Energy Mix  (donut chart from PNG)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Energy Mix", "Your energy — 100% renewable.")

png(s, "chart_donut", Inches(0.3), Inches(1.75), Inches(5.5), Inches(5.5))

callouts = [
    (LIGHT_GREEN, "60%", "On-site solar arrays"),
    (ACCENT_GOLD,  "30%", "Wind PPA — locked 15 yrs"),
    (MID_GREEN,    "10%", "Certified green grid backup"),
]
for i, (c, pct, desc) in enumerate(callouts):
    y = Inches(2.3 + i * 1.5)
    rect(s, Inches(6.3), y + Inches(0.1), Inches(0.08), Inches(0.7), c)
    txbox(s, pct, Inches(6.6), y, Inches(2), Inches(0.7),
          size=32, bold=True, colour=WHITE)
    txbox(s, desc, Inches(8.7), y + Inches(0.1), Inches(4.3), Inches(0.6),
          size=17, colour=LIGHT_GREY)
txbox(s, "100% renewable from day one — contractual guarantee.",
      Inches(6.3), Inches(6.4), Inches(6.5), Inches(0.5),
      size=14, bold=True, colour=LIGHT_GREEN)
label_tag(s, "Speaker: Person B")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Sustainability  (carbon line chart from PNG)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Sustainability", "Not just compliant — leading.")

png(s, "chart_carbon", Inches(0.3), Inches(1.9), Inches(7.3), Inches(4.5))

add_bullets(s, [
    "EU Taxonomy  ·  CSRD  ·  EN 50600",
    "Net-zero Y4 — operational, not offset",
    "Circular hardware lifecycle",
    "30% less embodied carbon",
], Inches(7.9), Inches(2.2), Inches(5.0), Inches(4.0), size=19)
label_tag(s, "Speaker: Person B  ·  4 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Risk Register  (table + risk matrix PNG)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Risk Register", "Four risks. All managed.")

styled_table(s,
    ["Risk", "Likelihood", "Impact", "Mitigation"],
    [
        ["Supply chain delay",   "Medium", "High",   "Pre-ordered long-lead items"],
        ["Energy price spike",   "Low",    "Medium", "15-yr wind PPA locked"],
        ["Regulatory change",    "Low",    "High",   "EU-aligned from day one"],
        ["Construction overrun", "Medium", "High",   "8% contingency built in"],
    ],
    Inches(0.4), Inches(2.0), Inches(7.8), Inches(4.1)
)
png(s, "chart_risk", Inches(8.4), Inches(1.9), Inches(4.7), Inches(4.5))
label_tag(s, "Speaker: Person B  ·  3 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Financial Overview  (bar chart PNG + KPI cards)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Financial Overview", "€50M. Fully allocated.")

png(s, "chart_finance", Inches(0.3), Inches(1.85), Inches(8.2), Inches(4.9))

kpis = [
    (LIGHT_GREEN, "€50M",   "Total investment"),
    (ACCENT_GOLD, "Year 6", "Breakeven"),
    (LIGHT_GREEN, "12%",    "IRR"),
]
for i, (c, big, small) in enumerate(kpis):
    y = Inches(2.0 + i * 1.55)
    rect(s, Inches(8.9), y, Inches(4.0), Inches(1.3), RGBColor(0x22, 0x4D, 0x3E))
    rect(s, Inches(8.9), y, Inches(0.10), Inches(1.3), c)
    txbox(s, big, Inches(9.2), y + Inches(0.05), Inches(3.6), Inches(0.75),
          size=34, bold=True, colour=c)
    txbox(s, small, Inches(9.2), y + Inches(0.78), Inches(3.6), Inches(0.45),
          size=14, colour=LIGHT_GREY)
label_tag(s, "Speaker: Person A  ·  4 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Roadmap  (Gantt PNG)
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Roadmap", "4 years. 4 phases.")
png(s, "chart_gantt", Inches(0.3), Inches(1.95), Inches(12.7), Inches(4.7))
label_tag(s, "Speaker: Person A  ·  2 min")

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Next Steps
# ════════════════════════════════════════════════════════════════════════════
s = add_slide(); bg(s)
slide_header(s, "Next Steps", "Three steps to begin.")

steps = [
    ("01", "Sign Letter\nof Intent", "by [date]"),
    ("02", "Detailed design\ndelivered", "Week 6"),
    ("03", "Kick-off\nmeeting", "Week 8"),
]
for i, (num, title, when) in enumerate(steps):
    x = Inches(0.5 + i * 4.2)
    rect(s, x, Inches(2.1), Inches(3.9), Inches(3.8), RGBColor(0x22, 0x4D, 0x3E))
    rect(s, x, Inches(2.1), Inches(3.9), Inches(0.12), MID_GREEN)
    txbox(s, num, x + Inches(0.15), Inches(2.3), Inches(3.6), Inches(0.8),
          size=40, bold=True, colour=MID_GREEN, align=PP_ALIGN.CENTER)
    txbox(s, title, x + Inches(0.15), Inches(3.1), Inches(3.6), Inches(1.1),
          size=20, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)
    txbox(s, when, x + Inches(0.15), Inches(4.2), Inches(3.6), Inches(0.6),
          size=16, colour=LIGHT_GREEN, align=PP_ALIGN.CENTER)

rect(s, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.75), RGBColor(0x1E, 0x44, 0x36))
txbox(s, "Your portal is live — login details on the card in front of you.",
      Inches(0.6), Inches(6.25), Inches(12), Inches(0.55),
      size=17, bold=True, colour=LIGHT_GREEN, align=PP_ALIGN.CENTER)
label_tag(s, "Speaker: Person A  ·  2 min")

# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "NexVerde_Pitch.pptx")
prs.save(out)
print(f"Saved: {out}")
